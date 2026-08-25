"""生成汇报 PPT：python tools/make_ppt.py
在脚本底部 STUDENT 处填入姓名/学号后重新运行即可更新封面。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path(__file__).resolve().parent.parent
SHOTS = BASE / "docs" / "screenshots"
OUT = BASE / "ppt" / "企业员工管理系统-项目汇报.pptx"

# ====== 在此修改答辩人信息 ======
STUDENT = {"name": "（姓名）", "student_id": "（学号）", "date": "2026-08-24"}

BLUE = RGBColor(0x1F, 0x3B, 0x73)
LIGHT = RGBColor(0x90, 0x99, 0x99)
DARK = RGBColor(0x30, 0x31, 0x33)


def new_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "微软雅黑"
    run.font.color.rgb = BLUE
    return slide


def add_text(slide, left, top, width, height, lines, size=16, color=DARK, bold=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "微软雅黑"
        run.font.color.rgb = color
    return tb


def add_pic(slide, name, left, top, width):
    path = SHOTS / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_textbox(Inches(0), Inches(2.0), Inches(13.333), Inches(3.0))
    tf = bg.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "企业员工管理系统"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.name = "微软雅黑"
    run.font.color.rgb = BLUE
    p.alignment = 1
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "—— 全栈 AI 辅助开发项目汇报 ——"
    run2.font.size = Pt(20)
    run2.font.name = "微软雅黑"
    run2.font.color.rgb = LIGHT
    p2.alignment = 1
    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = f"答辩人：{STUDENT['name']}    学号：{STUDENT['student_id']}\n{STUDENT['date']}"
    run3.font.size = Pt(16)
    run3.font.name = "微软雅黑"
    run3.font.color.rgb = LIGHT
    p3.alignment = 1


def build():
    prs = new_presentation()
    cover(prs)

    # 项目背景
    s = add_slide(prs, "一、项目背景与目标")
    add_text(s, 0.8, 1.4, 11.7, 5.5, [
        "某中小企业需要一套员工信息管理系统，管理：员工档案 / 部门架构 / 职位体系",
        "日常考勤（含打卡） / 请假审批 / 薪资核算 / 数据可视化仪表板",
        "",
        "目标：提升 HR 管理效率，规范员工行为记录，为管理层提供数据决策支持",
        "",
        "▸ 4 种角色：系统管理员 / HR / 部门主管 / 普通员工",
        "▸ 技术栈：Vue3 + FastAPI + SQLite（前后端分离）",
        "▸ 开发方式：AI 辅助编码（需求拆解 → 设计 → 编码 → 测试全流程）",
    ], size=18)

    # 技术架构
    s = add_slide(prs, "二、技术架构")
    add_text(s, 0.8, 1.3, 11.7, 5.5, [
        "前端层：Vue 3 + Vite + Element Plus + ECharts + Pinia",
        "          路由守卫双鉴权（菜单按角色动态渲染）",
        "",
        "接口层：RESTful API（/api/v1），JWT 无状态认证，12 小时过期",
        "          每个接口声明角色权限（require_roles 依赖注入）",
        "",
        "后端层：FastAPI + SQLAlchemy 2.0（参数化查询防注入）",
        "          密码 PBKDF2 加盐哈希（10 万次迭代）",
        "",
        "数据层：SQLite 文件数据库，9 张业务表，3 个唯一约束",
        "",
        "部署形态：本地 uvicorn(8000) + Vite(3000) 代理转发；兼容 CloudStudio 一键启动",
    ], size=17)

    # 功能总览
    s = add_slide(prs, "三、功能总览（10 大模块）")
    add_text(s, 0.8, 1.3, 11.7, 5.5, [
        "① 部门管理：列表 + 人员统计 + CRUD",
        "② 职位管理：所属部门 + 职级（初/中/高/资深）",
        "③ 员工管理：CRUD + 关键词/部门/状态筛选 + Excel 批量导入导出",
        "④ 考勤管理：上下班打卡 + 考勤规则配置 + 月度统计（迟到/早退/缺卡/出勤率）",
        "⑤ 请假管理：员工申请 → 主管审批（校验部门归属）→ 记录查询",
        "⑥ 薪资管理：HR 按员工+月份录入（实发自动计算），员工查看本人薪资",
        "⑦ 数据看板：员工总数、入离职、部门/职级分布、出勤率对比、请假统计、薪资趋势",
        "⑧ 用户管理：账号与员工档案关联、角色分配、密码重置",
        "⑨ 员工工作台：个人档案、今日考勤、本月统计、最近薪资、快捷操作",
        "⑩ 智能考勤异常提醒（创新功能）：连续迟到 / 缺卡 / 请假超额度 / 部门出勤率预警",
    ], size=16)

    # 演示页1：登录+看板
    s = add_slide(prs, "四、系统演示（1/3）登录与数据看板")
    add_pic(s, "01-登录页.png", 0.4, 1.2, 6.0)
    add_pic(s, "02-HR数据看板.png", 6.8, 1.2, 6.1)

    # 演示页2：员工管理+考勤
    s = add_slide(prs, "四、系统演示（2/3）员工管理与考勤")
    add_pic(s, "05-员工管理.png", 0.4, 1.2, 6.0)
    add_pic(s, "07-月度考勤统计.png", 6.8, 1.2, 6.1)

    # 演示页3：请假+薪资+用户
    s = add_slide(prs, "四、系统演示（3/3）请假审批与薪资")
    add_pic(s, "16-请假审批.png", 0.4, 1.2, 6.0)
    add_pic(s, "09-薪资管理.png", 6.8, 1.2, 6.1)
    add_pic(s, "10-用户管理.png", 4.6, 4.4, 4.2)

    # 创新点
    s = add_slide(prs, "五、创新点")
    add_text(s, 0.8, 1.3, 11.7, 5.5, [
        "1. 智能考勤异常提醒（规则引擎式自动化）",
        "   自动检测：连续 3 天迟到（严重）/ 当月缺卡 ≥3 次 / 请假超年度额度 / 部门出勤率 <90%",
        "   价值：替代 HR 人工翻查考勤记录，阈值配置化可调，可迭代接入 LLM 生成改善建议",
        "",
        "2. Excel 批量导入导出：导出模板即导入模板，逐行校验报错",
        "",
        "3. 权限体系：前端菜单动态渲染 + 后端接口级鉴权双重控制",
        "",
        "4. 数据驱动决策：看板 6 类图表让管理层一屏掌握人事全局",
    ], size=17)

    # 测试与质量
    s = add_slide(prs, "六、测试与质量保障")
    add_text(s, 0.8, 1.3, 11.7, 5.5, [
        "自动化测试：接口层（登录/权限/CRUD/统计/导入导出）全部通过",
        "浏览器端到端：Playwright 驱动真实浏览器走通「提交请假 → 主管审批 → 状态更新」4/4 通过",
        "界面冒烟测试：13/13 页面关键元素断言通过，16 张截图存档",
        "",
        "开发中发现并修复 4 个缺陷：",
        "  · 薪资趋势月份方向错误（负偏移逻辑）",
        "  · Excel 导出中文文件名编码 500",
        "  · 打卡接口响应结构序列化错误",
        "  · 路由顺序被路径参数吞掉（预防性修复）",
        "",
        "安全设计：PBKDF2 加盐哈希 / JWT / 参数化查询 / 接口级权限 / 主管跨部门审批拦截",
    ], size=17)

    # 总结展望
    s = add_slide(prs, "七、总结与展望")
    add_text(s, 0.8, 1.3, 11.7, 5.5, [
        "已交付：",
        "  · 完整可运行的全栈系统（后端 10 路由模块 + 前端 16 页面）",
        "  · 全套设计文档（Spec / 架构 / 数据库 / API / 审查报告 / 测试文档 / 操作手册）",
        "  · 演示数据一键初始化，4 个演示账号",
        "",
        "展望：",
        "  · 考勤时间改由服务端认定，防止手动补卡",
        "  · 请假天数按工作日计算，对接法定节假日",
        "  · 智能提醒接入大模型，自然语言生成考勤分析与改善建议",
        "  · 数据库升级 PostgreSQL，支持更大规模",
        "",
        "收获：完整实践了 AI 辅助的全栈开发流程——需求 → 设计 → 编码 → 测试 → 部署，",
        "理解了先设计后编码、用测试保障质量的重要性。",
    ], size=17)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"PPT 已生成: {OUT}")


if __name__ == "__main__":
    build()
